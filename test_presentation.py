#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Test script to verify the generated PowerPoint presentation
"""

from pptx import Presentation
import sys


def test_presentation(filepath='金陵购车节政企通补贴平台合作方案.pptx'):
    """Test the generated PowerPoint presentation"""
    
    print("Testing PowerPoint presentation...")
    print("="*60)
    
    try:
        # Load the presentation
        prs = Presentation(filepath)
        
        # Test 1: Check number of slides
        expected_slides = 11
        actual_slides = len(prs.slides)
        assert actual_slides == expected_slides, f"Expected {expected_slides} slides, got {actual_slides}"
        print(f"✓ Test 1 PASSED: Correct number of slides ({actual_slides})")
        
        # Test 2: Check slide dimensions
        expected_width = 9144000  # 10 inches in EMUs
        expected_height = 6858000  # 7.5 inches in EMUs
        assert prs.slide_width == expected_width, f"Unexpected slide width"
        assert prs.slide_height == expected_height, f"Unexpected slide height"
        print(f"✓ Test 2 PASSED: Correct slide dimensions")
        
        # Test 3: Check that each slide has text content
        slides_with_content = 0
        for i, slide in enumerate(prs.slides, 1):
            has_text = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    has_text = True
                    break
            if has_text:
                slides_with_content += 1
        
        assert slides_with_content == expected_slides, f"Only {slides_with_content} slides have content"
        print(f"✓ Test 3 PASSED: All slides have text content")
        
        # Test 4: Check for key content in specific slides
        key_contents = {
            0: "金陵购车节",  # Title slide
            1: "政银合作",     # Background slide
            2: "聚焦核心客群", # Target customers
            4: "政企通",       # Overall solution
            9: "政银合作新标杆", # Value summary
            10: "谢谢聆听"     # Thank you
        }
        
        for slide_idx, expected_text in key_contents.items():
            slide = prs.slides[slide_idx]
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text
            
            assert expected_text in slide_text, \
                f"Slide {slide_idx + 1} missing expected text '{expected_text}'"
        
        print(f"✓ Test 4 PASSED: Key content verified in slides")
        
        # Test 5: Verify Chinese character support
        chinese_chars_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Check if text contains Chinese characters
                    if any('\u4e00' <= char <= '\u9fff' for char in shape.text):
                        chinese_chars_found = True
                        break
            if chinese_chars_found:
                break
        
        assert chinese_chars_found, "No Chinese characters found in presentation"
        print(f"✓ Test 5 PASSED: Chinese character support verified")
        
        # Print detailed slide information
        print("\n" + "="*60)
        print("Slide Details:")
        print("="*60)
        
        slide_titles = []
        for i, slide in enumerate(prs.slides, 1):
            # Try to extract title
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()[:50]  # First 50 chars
                    break
            slide_titles.append(f"Slide {i}: {title}")
        
        for title in slide_titles:
            print(f"  {title}")
        
        print("\n" + "="*60)
        print("✓ ALL TESTS PASSED!")
        print("="*60)
        print(f"\nPresentation file: {filepath}")
        print(f"Total slides: {actual_slides}")
        print(f"File appears to be valid and ready for use.")
        
        return True
        
    except FileNotFoundError:
        print(f"✗ ERROR: File '{filepath}' not found")
        print("Please run generate_presentation.py first to create the file.")
        return False
    except AssertionError as e:
        print(f"✗ TEST FAILED: {str(e)}")
        return False
    except Exception as e:
        print(f"✗ UNEXPECTED ERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == '__main__':
    filepath = sys.argv[1] if len(sys.argv) > 1 else '金陵购车节政企通补贴平台合作方案.pptx'
    success = test_presentation(filepath)
    sys.exit(0 if success else 1)
