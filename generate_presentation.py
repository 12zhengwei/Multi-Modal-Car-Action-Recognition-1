#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint Generator for Jinling Car Purchase Subsidy Platform Project
Generates a complete business presentation for China Merchants Bank
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os


class SubsidyPlatformPresentation:
    """Generator for the Jinling Car Purchase Subsidy Platform presentation"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Color scheme - China Merchants Bank colors
        self.cmb_red = RGBColor(204, 0, 0)
        self.cmb_gold = RGBColor(255, 215, 0)
        self.dark_gray = RGBColor(64, 64, 64)
        self.light_gray = RGBColor(192, 192, 192)
        self.white = RGBColor(255, 255, 255)
        
    def add_title_slide(self):
        """第1页：封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = '金陵购车节"政企通"补贴平台合作方案'
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.cmb_red
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.7))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = '暨招商银行零售业务增长新机遇'
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = self.dark_gray
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Logo text (since we can't add actual images easily)
        logo_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
        logo_frame = logo_box.text_frame
        logo_frame.text = '招商银行 + 金陵市商务局'
        logo_para = logo_frame.paragraphs[0]
        logo_para.font.size = Pt(20)
        logo_para.font.color.rgb = self.dark_gray
        logo_para.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.4))
        date_frame = date_box.text_frame
        date_frame.text = '2025年10月'
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = self.light_gray
        date_para.alignment = PP_ALIGN.CENTER
        
    def add_background_slide(self):
        """第2页：项目背景与机遇分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '政银合作，共促消费新篇章'
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        sections = [
            ("政府目标", "金陵市商务局致力于促进汽车消费、拉动内需，推动本地经济高质量发展"),
            ("项目概述", "活动时间：2025年第四季度\n补贴总额：2000万元\n补贴标准：3000-5000元/车"),
            ("我行角色", "作为独家合作银行，负责平台开发、资金发放与核销的核心工作"),
            ("机遇洞察", "这不仅是完成政府任务，更是品牌提升、客户增长、业务拓展的宝贵机遇")
        ]
        
        for i, (heading, content) in enumerate(sections):
            p = tf.add_paragraph()
            p.text = f"• {heading}："
            p.level = 0
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(6)
            
            p = tf.add_paragraph()
            p.text = content
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(12)
    
    def add_target_customers_slide(self):
        """第3页：目标客群精准定位"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '聚焦核心客群，最大化合作价值'
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = "核心目标人群"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "金陵市内购买新车的个人消费者"
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "两大核心客群画像"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(8)
        
        customer_profiles = [
            ("首购族青年", "对价格敏感，习惯线上操作，是潜在的信用卡和信贷用户"),
            ("家庭升级用户", "追求品质，有更高消费能力，是财富管理和综合金融服务的潜在客户")
        ]
        
        for heading, desc in customer_profiles:
            p = tf.add_paragraph()
            p.text = f"✓ {heading}："
            p.level = 1
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            p.text = desc
            p.level = 2
            p.font.size = Pt(15)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = "策略"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "通过精准定位，为不同客群设计针对性服务，实现价值最大化"
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = self.dark_gray
    
    def add_case_studies_slide(self):
        """第4页：他山之石：成功案例借鉴"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '借鉴卓越实践，启发我行方案'
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        insights = [
            ("启示一（建行案例）", "\"政府搭台、银行助力\"模式可快速起量，我行App可作为核心载体"),
            ("启示二（多银行实践）", "\"线上领券+线下核销\"模式成熟有效，是实现全场景消费闭环的关键"),
            ("启示三（邮储银行案例）", "\"政府补贴+银行优惠\"双重激励能极大激发热情，并有效带动信用卡和贷款业务")
        ]
        
        for heading, content in insights:
            p = tf.add_paragraph()
            p.text = f"• {heading}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(6)
            
            p = tf.add_paragraph()
            p.text = content
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(14)
        
        p = tf.add_paragraph()
        p.text = "结论"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.cmb_red
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "我们的方案必须超越一个简单的申请工具，构建一个\"获客-活客-留客\"的业务闭环"
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = self.dark_gray
    
    def add_overall_solution_slide(self):
        """第5页：总体方案设计：一站式"政企通"平台"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '一站式"政企通"补贴平台总体架构'
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = "平台三大核心组成部分"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(12)
        
        components = [
            ("C端用户入口", "\"补贴申请\"小程序（内嵌于招行App或微信生态）", "便捷的用户申请体验"),
            ("B端管理后台", "补贴申请审核平台（供银行和商务局人员使用）", "高效的审核管理工具"),
            ("底层支撑系统", "账户、支付、信审、数据分析等银行核心系统", "稳定可靠的技术基础")
        ]
        
        for title_text, desc, feature in components:
            p = tf.add_paragraph()
            p.text = f"▶ {title_text}"
            p.level = 0
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            p.text = desc
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(2)
            
            p = tf.add_paragraph()
            p.text = f"→ {feature}"
            p.level = 2
            p.font.size = Pt(14)
            p.font.color.rgb = self.light_gray
            p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = "设计理念：便捷 · 合规 · 智能 · 增长"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.cmb_red
        p.alignment = PP_ALIGN.CENTER
    
    def add_user_flow_slide(self):
        """第6页：核心流程设计：用户申请流程"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '5分钟极速申请：合规与便捷的完美平衡'
        title.text_frame.paragraphs[0].font.size = Pt(34)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(0.8)
        top = Inches(2)
        width = Inches(8.4)
        height = Inches(4.8)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        steps = [
            ("步骤1：入口/授权", "从招行App/小程序进入，一键授权登录"),
            ("步骤2：车辆信息录入", "OCR识别行驶证/购车发票，用户确认"),
            ("步骤3：证件验证", "OCR识别人脸/身份证，与公安系统联网核验"),
            ("步骤4：绑定收款卡", 
             "• 智能引导：优先推荐绑定招行借记卡\n"
             "• 若用户选择他行卡，弹窗提示\"使用招行卡可享额外好礼/更快到账\"\n"
             "• 引导一键开立电子账户\n"
             "• 系统自动校验银行卡有效性"),
            ("步骤5：提交申请与进度查询", 
             "• 提交后生成申请编号\n"
             "• 用户可随时在\"我的申请\"中查看审核状态\n"
             "• 状态：待审核/审核通过/已驳回/补贴已发放")
        ]
        
        for i, (step_title, step_desc) in enumerate(steps):
            p = tf.add_paragraph()
            p.text = step_title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            p.text = step_desc
            p.level = 1
            p.font.size = Pt(13)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(8)
    
    def add_c_end_design_slide(self):
        """第7页：C端核心：补贴申请小程序设计"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '打造用户体验极致的申请小程序'
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        features = [
            ("申请入口", "简洁明了的申请按钮，突出\"政府补贴\"字样，吸引用户关注"),
            ("智能填报", "借助OCR和大数据技术，自动识别和预填信息，减少用户操作，提升体验"),
            ("资格初审", "基于规则引擎，对用户填写的资料进行实时自动初审\n对明显不符的申请进行即时提示，减少无效申请流入审核后台"),
            ("二次申请通道", "对于被驳回的申请，清晰告知驳回原因\n提供\"补充材料/修改申请\"的便捷通道，提升用户体验")
        ]
        
        for feature_title, feature_desc in features:
            p = tf.add_paragraph()
            p.text = f"◆ {feature_title}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(6)
            
            p = tf.add_paragraph()
            p.text = feature_desc
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(12)
    
    def add_b_end_design_slide(self):
        """第8页：B端核心：审核平台设计"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '高效智能的审核平台，赋能管理'
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        features = [
            ("审核工作台", "清晰展示待审核、已通过、已驳回的申请列表\n支持多维度筛选和排序，提升审核效率"),
            ("详情审查页", "集中展示单笔申请的所有信息（申请人、车辆、材料照片）\n关键信息高亮显示，便于快速审查"),
            ("一键审核", "审核员可一键\"通过\"或\"驳回\"\n驳回时需选择或填写原因，确保流程规范"),
            ("数据驾驶舱", "实时展示申请总量、审核进度、补贴发放总额、客群分析等核心KPI\n为项目决策提供数据支持，助力精准运营")
        ]
        
        for feature_title, feature_desc in features:
            p = tf.add_paragraph()
            p.text = f"◆ {feature_title}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(6)
            
            p = tf.add_paragraph()
            p.text = feature_desc
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(10)
    
    def add_growth_strategy_slide(self):
        """第9页：业务增长策略：超越补贴，实现多赢"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '不止于补贴：构建零售业务增长新引擎'
        title.text_frame.paragraphs[0].font.size = Pt(34)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(0.8)
        top = Inches(2)
        width = Inches(8.4)
        height = Inches(4.8)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = "一、精准引流信贷业务"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "场景一（申请时）："
        p.level = 1
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = "提交申请后，基于用户资质（如车辆价格）和数据模型，精准推送\"招行车主信用卡\"或\"个人消费贷款\"推荐，额度预估，一键申请"
        p.level = 2
        p.font.size = Pt(14)
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = "文案：\"恭喜！您已获得XX万元招行备用金资格，购车消费更轻松！\""
        p.level = 2
        p.font.size = Pt(13)
        p.font.color.rgb = self.light_gray
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "场景二（等待时）："
        p.level = 1
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = "在申请进度查询页面，植入我行金融产品广告"
        p.level = 2
        p.font.size = Pt(14)
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = "二、提升平台用户活跃度"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(8)
        
        strategies = [
            ("引导绑卡", "强力引导用户使用招行借记卡领补贴，将其转化为我行有效客户"),
            ("App引流", "整个流程设计优先在招行App内完成，补贴发放后，通过App消息推送引导用户参与其他活动（如本地生活优惠、理财产品推荐），提升App月活"),
            ("业务闭环", "实现\"获客-活客-留客\"闭环：通过本次活动获取新客，通过App互动激活客户，通过后续的综合金融服务留住客户")
        ]
        
        for strategy_title, strategy_desc in strategies:
            p = tf.add_paragraph()
            p.text = f"• {strategy_title}："
            p.level = 1
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            p.text = strategy_desc
            p.level = 2
            p.font.size = Pt(13)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(6)
    
    def add_value_summary_slide(self):
        """第10页：项目价值与总结"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = '政银合作新标杆，共创社会与商业价值'
        title.text_frame.paragraphs[0].font.size = Pt(34)
        title.text_frame.paragraphs[0].font.color.rgb = self.cmb_red
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        values = [
            ("对政府", "高效、透明、精准地落实了民生补贴政策\n提升了政府公共服务的现代化水平"),
            ("对市民", "提供了便捷、优质的补贴申请体验\n真正享受到政府惠民政策的实惠"),
            ("对招行", "")
        ]
        
        for stakeholder, value in values:
            p = tf.add_paragraph()
            p.text = f"▶ {stakeholder}"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(6)
            
            if value:
                p = tf.add_paragraph()
                p.text = value
                p.level = 1
                p.font.size = Pt(16)
                p.font.color.rgb = self.dark_gray
                p.space_after = Pt(12)
        
        # Special formatting for 招行 section
        bank_values = [
            ("品牌价值", "树立了有担当、有能力的社会责任形象"),
            ("商业价值", "批量获取了高质量车主客户，并成功带动了信用卡、贷款、App活跃度等多项业务指标的增长")
        ]
        
        for value_type, value_desc in bank_values:
            p = tf.add_paragraph()
            p.text = f"  • {value_type}："
            p.level = 1
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(4)
            
            p = tf.add_paragraph()
            p.text = value_desc
            p.level = 2
            p.font.size = Pt(15)
            p.font.color.rgb = self.dark_gray
            p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = "恳请领导审议，我们有信心将此项目打造成政银合作的典范！"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.cmb_red
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)
    
    def add_thank_you_slide(self):
        """第11页：封底页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Thank you message
        thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        thank_frame = thank_box.text_frame
        thank_frame.text = '谢谢聆听'
        thank_para = thank_frame.paragraphs[0]
        thank_para.font.size = Pt(48)
        thank_para.font.bold = True
        thank_para.font.color.rgb = self.cmb_red
        thank_para.alignment = PP_ALIGN.CENTER
        
        # Contact info
        contact_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.6))
        contact_frame = contact_box.text_frame
        contact_frame.text = '招商银行零售金融部'
        contact_para = contact_frame.paragraphs[0]
        contact_para.font.size = Pt(20)
        contact_para.font.color.rgb = self.dark_gray
        contact_para.alignment = PP_ALIGN.CENTER
        
        # Logo text
        logo_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
        logo_frame = logo_box.text_frame
        logo_frame.text = '招商银行'
        logo_para = logo_frame.paragraphs[0]
        logo_para.font.size = Pt(24)
        logo_para.font.color.rgb = self.cmb_red
        logo_para.alignment = PP_ALIGN.CENTER
    
    def generate(self, output_path='金陵购车节政企通补贴平台合作方案.pptx'):
        """Generate the complete presentation"""
        print("Generating PowerPoint presentation...")
        
        # Add all slides
        self.add_title_slide()
        print("✓ Added slide 1: Cover page")
        
        self.add_background_slide()
        print("✓ Added slide 2: Project background")
        
        self.add_target_customers_slide()
        print("✓ Added slide 3: Target customers")
        
        self.add_case_studies_slide()
        print("✓ Added slide 4: Case studies")
        
        self.add_overall_solution_slide()
        print("✓ Added slide 5: Overall solution")
        
        self.add_user_flow_slide()
        print("✓ Added slide 6: User flow")
        
        self.add_c_end_design_slide()
        print("✓ Added slide 7: C-end design")
        
        self.add_b_end_design_slide()
        print("✓ Added slide 8: B-end design")
        
        self.add_growth_strategy_slide()
        print("✓ Added slide 9: Growth strategy")
        
        self.add_value_summary_slide()
        print("✓ Added slide 10: Value summary")
        
        self.add_thank_you_slide()
        print("✓ Added slide 11: Thank you page")
        
        # Save presentation
        self.prs.save(output_path)
        print(f"\n✓ Presentation saved to: {output_path}")
        print(f"Total slides: {len(self.prs.slides)}")
        
        return output_path


def main():
    """Main function to generate the presentation"""
    generator = SubsidyPlatformPresentation()
    output_file = generator.generate()
    
    print("\n" + "="*60)
    print("PowerPoint presentation generated successfully!")
    print("="*60)
    print(f"\nFile: {output_file}")
    print("\nThe presentation includes 11 slides:")
    print("1. Cover page")
    print("2. Project background and opportunities")
    print("3. Target customer positioning")
    print("4. Case study references")
    print("5. Overall solution design")
    print("6. Core user flow design")
    print("7. C-end application design")
    print("8. B-end management platform design")
    print("9. Business growth strategy")
    print("10. Project value summary")
    print("11. Thank you page")
    print("\nYou can now open the .pptx file in PowerPoint or compatible software.")


if __name__ == '__main__':
    main()
